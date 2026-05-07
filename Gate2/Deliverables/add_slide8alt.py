"""
Add Slide 8 (Alt) — Non-Determinism × Volume Grid — after slide 8 in the existing PPTX.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import copy

# ── Colour palette (matches build_pptx.py) ─────────────────────────────────
APEX_DARK  = RGBColor(0x1A, 0x2E, 0x44)
APEX_MID   = RGBColor(0x1F, 0x61, 0xAA)
APEX_LIGHT = RGBColor(0xE8, 0xF1, 0xFB)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY_TEXT  = RGBColor(0x44, 0x44, 0x44)
GREEN_COL  = RGBColor(0x00, 0x7A, 0x33)
AMBER      = RGBColor(0xFF, 0x99, 0x00)
RED_COL    = RGBColor(0xCC, 0x00, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

PPTX_PATH = r"c:\Users\Benoit_Charrier\FDE Program\FDE-Program\Gate2\Deliverables\Stakeholder_Presentation.pptx"

prs = Presentation(PPTX_PATH)
blank_layout = prs.slide_layouts[6]


def set_bg(slide, colour):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, colour=GREY_TEXT,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return txBox


def add_header_bar(slide, title_text, subtitle_text=""):
    bar = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0), SLIDE_W, Inches(1.4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = APEX_DARK
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle_text
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)
    tf.margin_left = Inches(0.3)
    tf.margin_top  = Inches(0.18)


def add_footer(slide, text="Apex Distribution Ltd  |  Billing Disputes — Assessment & Proposed Solution  |  Confidential"):
    add_text_box(slide, text,
                 Inches(0.2), Inches(7.15), Inches(12.9), Inches(0.3),
                 font_size=9, colour=RGBColor(0x99, 0x99, 0x99),
                 align=PP_ALIGN.CENTER)


def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


# ── Build Slide 8 (Alt) ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)   # added at the end for now
set_bg(s, WHITE)
add_header_bar(s, "The Opportunity — Non-Determinism × Volume Grid",
               "Alternative to Slide 8 (choose one before presenting)")

# ── Quadrant chart ────────────────────────────────────────────────────────────
# Chart area
chart_left = Inches(0.55)
chart_top  = Inches(1.62)
chart_w    = Inches(7.4)
chart_h    = Inches(5.3)

# Background: lightly shade the Q1 zone (top-right = primary agentic targets)
q1_bg = s.shapes.add_shape(
    1,
    chart_left + chart_w / 2, chart_top,
    chart_w / 2, chart_h / 2
)
q1_bg.fill.solid()
q1_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xEA)
q1_bg.line.fill.background()

# Q2 bg (top-left)
q2_bg = s.shapes.add_shape(1, chart_left, chart_top, chart_w / 2, chart_h / 2)
q2_bg.fill.solid()
q2_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE8)
q2_bg.line.fill.background()

# Q3 bg (bottom-left)
q3_bg = s.shapes.add_shape(1, chart_left, chart_top + chart_h / 2, chart_w / 2, chart_h / 2)
q3_bg.fill.solid()
q3_bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
q3_bg.line.fill.background()

# Q4 bg (bottom-right)
q4_bg = s.shapes.add_shape(1, chart_left + chart_w / 2, chart_top + chart_h / 2,
                             chart_w / 2, chart_h / 2)
q4_bg.fill.solid()
q4_bg.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
q4_bg.line.fill.background()

# Axis lines
ax_h = s.shapes.add_shape(1, chart_left, chart_top + chart_h / 2 - Inches(0.02),
                            chart_w, Inches(0.04))
ax_h.fill.solid(); ax_h.fill.fore_color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
ax_h.line.fill.background()

ax_v = s.shapes.add_shape(1, chart_left + chart_w / 2 - Inches(0.02), chart_top,
                            Inches(0.04), chart_h)
ax_v.fill.solid(); ax_v.fill.fore_color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
ax_v.line.fill.background()

# Quadrant labels (small, muted)
quad_labels = [
    ("Primary agentic targets",   chart_left + chart_w * 0.53,  chart_top + Inches(0.08),  GREEN_COL),
    ("Rules / RPA only",          chart_left + Inches(0.1),      chart_top + Inches(0.08),  AMBER),
    ("Not worth automating",      chart_left + Inches(0.1),      chart_top + chart_h * 0.55, GREY_TEXT),
    ("Select agentic use cases",  chart_left + chart_w * 0.53,  chart_top + chart_h * 0.55, APEX_MID),
]
for ql_text, ql_x, ql_y, ql_col in quad_labels:
    add_text_box(s, ql_text, ql_x, ql_y, Inches(3.3), Inches(0.28),
                 font_size=9, colour=ql_col, bold=True)

# Axis titles
add_text_box(s, "← Low Non-Determinism                    High Non-Determinism →",
             chart_left, chart_top + chart_h + Inches(0.08), chart_w, Inches(0.28),
             font_size=10, colour=GREY_TEXT, align=PP_ALIGN.CENTER)
add_text_box(s, "Low\nVolume",
             chart_left - Inches(0.5), chart_top + chart_h - Inches(0.7), Inches(0.45), Inches(0.6),
             font_size=9, colour=GREY_TEXT, align=PP_ALIGN.CENTER)
add_text_box(s, "High\nVolume",
             chart_left - Inches(0.5), chart_top, Inches(0.45), Inches(0.6),
             font_size=9, colour=GREY_TEXT, align=PP_ALIGN.CENTER)

# Work-stream dots
# Mermaid coords [x=non-determinism 0→1, y=volume 0→1]
# In chart pixel space: x grows left→right, y grows bottom→top (invert y)
ws_data = [
    # name,  nd,   vol,   colour,   diam,     primary,  note
    ("WS4\nBilling Disputes\n★ PRIMARY",   0.92, 0.75, APEX_MID,   Inches(1.05), True,  ""),
    ("WS1\nDelivery Exceptions",           0.75, 0.75, GREY_TEXT,  Inches(0.85), False, ""),
    ("WS3\nDispatch Adjustments*",         0.67, 0.65, GREY_TEXT,  Inches(0.85), False, "*excluded"),
    ("WS2\nETA Inquiries",                 0.25, 0.92, GREY_TEXT,  Inches(0.85), False, ""),
]

for (name, nd, vol, col, diam, primary, note) in ws_data:
    cx = chart_left + chart_w  * nd  - diam / 2
    cy = chart_top  + chart_h  * (1 - vol) - diam / 2   # invert y
    dot = s.shapes.add_shape(9, cx, cy, diam, diam)
    dot.fill.solid()
    dot.fill.fore_color.rgb = APEX_MID if primary else RGBColor(0x88, 0xAA, 0xCC)
    dot.line.fill.background()
    # Label below dot
    label_y = cy + diam + Inches(0.04)
    add_text_box(s, name, cx - Inches(0.1), label_y,
                 diam + Inches(0.2), Inches(0.72),
                 font_size=9 if not primary else 10,
                 bold=primary,
                 colour=APEX_DARK if primary else GREY_TEXT,
                 align=PP_ALIGN.CENTER)

# ── Metrics panel on the right ────────────────────────────────────────────────
panel = s.shapes.add_shape(1, Inches(8.3), Inches(1.62), Inches(4.7), Inches(5.3))
panel.fill.solid(); panel.fill.fore_color.rgb = APEX_LIGHT
panel.line.color.rgb = APEX_MID

metrics = [
    ("Agentic Value Score",  "20 / 25  — highest reasoning complexity\n(Non-Determinism 5) combined with confirmed daily volume"),
    ("Annual baseline cost", "~£245k / year at current handle time"),
    ("Projected agent cost", "~£70k / year (incl. human review time)"),
    ("Directional saving",   "~£175k / year"),
    ("Build cost estimate",  "~£100k"),
    ("Payback period",       "~7 months"),
]
my = Inches(1.82)
add_text_box(s, "WS4 — Billing Disputes", Inches(8.5), my, Inches(4.3), Inches(0.4),
             font_size=13, bold=True, colour=APEX_DARK)
my += Inches(0.5)
for label, val in metrics:
    add_text_box(s, label, Inches(8.5), my, Inches(4.3), Inches(0.28),
                 font_size=11, bold=True, colour=APEX_DARK)
    my += Inches(0.3)
    add_text_box(s, val, Inches(8.5), my, Inches(4.3), Inches(0.42),
                 font_size=11, colour=APEX_MID)
    my += Inches(0.5)

add_text_box(s, "*WS3 plots in Q1 by score but is excluded — dispatch console has no confirmed programmatic interface (D3 §1)",
             Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.28),
             font_size=9, colour=GREY_TEXT)

add_footer(s)
add_speaker_notes(s,
    "The axes here are non-determinism — how much reasoning, judgment, and synthesis the work "
    "actually requires — and volume. An agent earns its keep in the top-right quadrant: high enough "
    "volume that the time saving compounds, and high enough reasoning complexity that a simpler "
    "automation can't do the job. ETA inquiries score high on volume but low on non-determinism — "
    "that's a lookup, not an agent. Dispatch adjustments would sit in the primary target zone if "
    "the dispatch console had a confirmed API surface; it doesn't yet, so it stays conditional. "
    "Billing disputes is the target: highest non-determinism in the portfolio, confirmed data "
    "access, and an active compliance gap the agent closes at the same time as it cuts handle time.")


# ── Move the new slide from the end to position 8 (after current slide 8) ────
# python-pptx stores slides in prs.slides._sldIdLst; we reorder the XML element.
xml_slides = prs.slides._sldIdLst
# The new slide is currently at index -1 (last). Move it to index 8 (0-based).
new_slide_el = xml_slides[-1]
xml_slides.remove(new_slide_el)
xml_slides.insert(8, new_slide_el)

prs.save(PPTX_PATH)
print(f"Slide 8 (Alt) inserted at position 9.  Saved: {PPTX_PATH}")
