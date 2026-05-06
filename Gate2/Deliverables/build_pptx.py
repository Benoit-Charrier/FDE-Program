"""
Convert Stakeholder_Presentation.md to a PowerPoint file.
Run from the Gate2/Deliverables directory or anywhere with the .venv active.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
APEX_DARK  = RGBColor(0x1A, 0x2E, 0x44)   # deep navy
APEX_MID   = RGBColor(0x1F, 0x61, 0xAA)   # blue
APEX_LIGHT = RGBColor(0xE8, 0xF1, 0xFB)   # pale blue (table shading)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY_TEXT  = RGBColor(0x44, 0x44, 0x44)
AMBER      = RGBColor(0xFF, 0x99, 0x00)
RED_COL    = RGBColor(0xCC, 0x00, 0x00)
GREEN_COL  = RGBColor(0x00, 0x7A, 0x33)

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, colour):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, colour=GREY_TEXT,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return txBox


def add_header_bar(slide, title_text, subtitle_text=""):
    """Dark navy header bar across the top."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_W, Inches(1.4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = APEX_DARK
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle_text:
        from pptx.oxml.ns import qn
        from lxml import etree
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle_text
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)

    # reposition text inside bar
    tf.margin_left  = Inches(0.3)
    tf.margin_top   = Inches(0.18)


def add_footer(slide, text="Apex Distribution Ltd  |  Billing Disputes — Assessment & Proposed Solution  |  Confidential"):
    add_text_box(slide, text,
                 Inches(0.2), Inches(7.15), Inches(12.9), Inches(0.3),
                 font_size=9, colour=RGBColor(0x99, 0x99, 0x99),
                 align=PP_ALIGN.CENTER)


def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


# ────────────────────────────────────────────────────────────────────────────
# Build presentation
# ────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # truly blank


# ── SLIDE 1 — Title ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)

# Full-height navy block on left third
left_bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(4.5), SLIDE_H)
left_bar.fill.solid()
left_bar.fill.fore_color.rgb = APEX_DARK
left_bar.line.fill.background()

add_text_box(s, "Apex Distribution Ltd",
             Inches(0.3), Inches(1.8), Inches(3.9), Inches(0.7),
             font_size=22, bold=True, colour=WHITE)
add_text_box(s, "Customer Operations — Billing Disputes",
             Inches(0.3), Inches(2.55), Inches(3.9), Inches(0.5),
             font_size=14, colour=RGBColor(0xAA, 0xCC, 0xEE))

add_text_box(s, "Assessment Findings\n& Proposed Solution",
             Inches(5.0), Inches(2.2), Inches(7.8), Inches(1.5),
             font_size=34, bold=True, colour=APEX_DARK)
add_text_box(s, "Date: [DATE]\nPresenter: [NAME], Field Deployment Engineer",
             Inches(5.0), Inches(4.0), Inches(7.8), Inches(0.9),
             font_size=14, colour=GREY_TEXT)
add_footer(s)
add_speaker_notes(s,
    "Thank you for making time for this. Today I'll walk you through what we found across "
    "your billing dispute process, what we think is worth building, and — critically — where "
    "we need your input before any build begins. This is not a pitch for AI. It's a structured "
    "look at where a well-scoped agent would actually save time and close a compliance gap that "
    "already exists in your process today, followed by an honest conversation about the decisions "
    "you need to make first.")


# ── SLIDE 2 — Agenda ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "Agenda")

items = [
    "1.  The business problem — what triggered this assessment",
    "2.  What we found — how the work flows today and where it gets stuck",
    "3.  What we recommend — a targeted agent for billing disputes",
    "4.  Decisions we need from you — five questions that determine the design",
    "5.  Next steps — four actions before build begins",
]
y = Inches(1.65)
for item in items:
    add_text_box(s, item, Inches(0.6), y, Inches(12.1), Inches(0.45),
                 font_size=18, colour=GREY_TEXT)
    y += Inches(0.78)

add_footer(s)
add_speaker_notes(s,
    "We'll move through five sections. The first two are about what we observed. The third is "
    "our recommendation. The fourth is where I'll need your answers — not hypotheticals, but "
    "specific operational facts that change what we build. The fifth is a concrete action list. "
    "I'll leave time for open discussion before next steps.")


# ── SLIDE 3 — Four Work Streams ───────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "Four Work Streams — Why We Focused on Billing Disputes")

# Table data
headers = ["Work stream", "Daily volume", "Time per case", "Assessment outcome"]
rows = [
    ["ETA inquiries",       "~400/day", "4 min",  "Automation candidate — primarily a lookup task; lower complexity"],
    ["Delivery exceptions", "~180/day", "12 min", "Future candidate — Driver App integration unconfirmed"],
    ["Dispatch adjustments","~90/day",  "18 min", "Not yet buildable — dispatch console has no confirmed API"],
    ["Billing disputes ★",  "~60/day",  "28 min", "PRIMARY TARGET — highest handle time, compliance gap, confirmed data access"],
]

col_widths = [Inches(2.4), Inches(1.4), Inches(1.4), Inches(7.4)]
table_left = Inches(0.3)
table_top  = Inches(1.55)
table_w    = sum(col_widths)
table_h    = Inches(2.6)

tbl = s.shapes.add_table(5, 4, table_left, table_top, table_w, table_h).table

for ci, (w, h) in enumerate(zip(col_widths, headers)):
    tbl.columns[ci].width = w
    cell = tbl.cell(0, ci)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = APEX_DARK
    p = cell.text_frame.paragraphs[0]
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE
    p.runs[0].font.size = Pt(12)

for ri, row_data in enumerate(rows, start=1):
    highlight = (ri == 4)
    for ci, val in enumerate(row_data):
        cell = tbl.cell(ri, ci)
        cell.text = val
        if highlight:
            cell.fill.solid()
            cell.fill.fore_color.rgb = APEX_LIGHT
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(11)
        p.runs[0].font.bold = highlight
        p.runs[0].font.color.rgb = APEX_DARK if highlight else GREY_TEXT

# Rationale bullets
bullets = [
    "Dispatch adjustments: Citrix-hosted system, no confirmed API — same failure mode as 2024 RPA",
    "ETA inquiries: high volume, low complexity — automation, not an agent (Agentic Value Score 10)",
    "Delivery exceptions: strong future target once Driver App integration is confirmed (score 16, conditional)",
]
y = Inches(4.35)
for b in bullets:
    add_text_box(s, "▸  " + b, Inches(0.4), y, Inches(12.5), Inches(0.38),
                 font_size=12, colour=GREY_TEXT)
    y += Inches(0.42)

add_footer(s)
add_speaker_notes(s,
    "Before I show you what we found in billing disputes specifically, I want to be transparent "
    "about how we chose to focus there. We looked at all four work streams and applied a consistent "
    "assessment lens — volume, complexity, system readiness, and compliance risk. Dispatch adjustments "
    "score highly on complexity but the dispatch console is a Citrix-hosted Java application with no "
    "confirmed API surface, and building a brittle integration there would recreate exactly the failure "
    "mode of the 2024 RPA. ETA inquiries are your highest-volume work stream but they're mostly a "
    "lookup — that's an automation, not an agent. Billing disputes is where complexity, compliance risk, "
    "and confirmed data access converge. That's why it's the recommendation.")


# ── SLIDE 4 — Business Problem ────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "Why We Are Here — The Business Problem")

bullets = [
    ("60 billing disputes per day at 28 min each",
     "Highest handle time in Customer Operations — absorbing 1,600+ minutes of skilled agent time daily with no structured support"),
    ("Credits applied without an audit trail",
     "Live case records show at least one £170 credit applied with no named approver and no audit log entry (Artefact 2, D3)"),
    ("Competitor benchmark: £1.2M annualised saving",
     "CEO has asked whether Apex can achieve something comparable (scenario_context.md)"),
]

y = Inches(1.65)
for bold_part, detail in bullets:
    # Bold heading
    add_text_box(s, bold_part, Inches(0.6), y, Inches(12.1), Inches(0.38),
                 font_size=16, bold=True, colour=APEX_DARK)
    y += Inches(0.38)
    add_text_box(s, detail, Inches(0.9), y, Inches(11.8), Inches(0.42),
                 font_size=14, colour=GREY_TEXT)
    y += Inches(0.55)

# Key question box
q_box = s.shapes.add_shape(1, Inches(0.4), Inches(5.6), Inches(12.2), Inches(0.9))
q_box.fill.solid()
q_box.fill.fore_color.rgb = APEX_LIGHT
q_box.line.color.rgb = APEX_MID

add_text_box(s,
    "Assessment question: Can a targeted agent reduce dispute handling time, close the compliance gap, "
    "and give Apex an auditable trace of every credit decision — without repeating the failures of 2024?",
    Inches(0.55), Inches(5.65), Inches(12.0), Inches(0.82),
    font_size=13, colour=APEX_DARK)

add_footer(s)
add_speaker_notes(s,
    "The 2024 chatbot and the billing RPA are in the room. I want to address them directly: both failed "
    "for specific, diagnosable reasons — the chatbot was customer-facing without a clear job to do, and "
    "the RPA broke because Aurum's data format changed without warning. What we're recommending today is "
    "different in both scope and design, and I'll show you specifically what we've built in to prevent "
    "the Aurum schema failure from happening again. But first, let me show you what the process actually "
    "looks like today.")


# ── SLIDE 5 — Process Flow ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "How the Work Actually Flows Today")

flow_steps = [
    "Customer emails billing@ with disputed charge",
    "Agent opens CRM → manually retrieves Aurum invoice data (avg 9 days to resolution — Artefact 2)",
    "Agent assesses whether the charge is valid — no formal policy to check against",
    "Agent decides credit amount from experience (e.g. £170 on £340 dispute) — Sandra's heuristic, not an approved rule",
    "Agent applies credit via manual override — no audit log entry recorded",
    "Customer receives credit on next statement — case closed informally",
]

arrow = "↓"
x_left = Inches(0.5)
box_w  = Inches(12.2)
box_h  = Inches(0.55)
y = Inches(1.55)

for i, step in enumerate(flow_steps):
    bg_col = APEX_LIGHT if i % 2 == 0 else WHITE
    box = s.shapes.add_shape(1, x_left, y, box_w, box_h)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_col
    box.line.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    add_text_box(s, step, x_left + Inches(0.12), y + Inches(0.06),
                 box_w - Inches(0.25), box_h - Inches(0.1),
                 font_size=13, colour=APEX_DARK)
    y += box_h
    if i < len(flow_steps) - 1:
        add_text_box(s, arrow, Inches(6.4), y, Inches(0.5), Inches(0.25),
                     font_size=14, colour=APEX_MID, align=PP_ALIGN.CENTER)
        y += Inches(0.25)

add_footer(s)
add_speaker_notes(s,
    "The flow I've drawn is based on actual case records, not the SOP. The SOP references DispatchHub, "
    "which was retired eighteen months ago, and the section covering damaged consignments is blank — "
    "marked TBD. So your team is resolving sixty cases a day against no documented policy and using a "
    "system that wasn't designed to give them the information they need at the time they need it. The "
    "result is what you see in that internal note on the Hayes & Sons case: a credit applied with no "
    "record of who approved it.")


# ── SLIDE 6 — Cognitive Hotspots ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "Where Time Goes — The Cognitive Hotspots")

hotspots = [
    ("1. Assembling the evidence  (est. 10–12 min per case)",
     "Sandra manually retrieves invoice data, surcharge details, and delivery records from separate "
     "systems — Aurum batch exports and the CRM — that do not talk to each other in real time. "
     "There is no single view of a dispute. Each case starts from scratch.  (D1 — WS4 Zone 1, Z2)"),
    ("2. Deciding the credit amount  (judgment call with no rule)",
     "No written policy exists for how much to credit. The observed practice — roughly 50% of the "
     "disputed amount — is Sandra's heuristic, not an approved rule. Different agents would reach "
     "different amounts for the same case.  (D1 — WS4 Zone 3, MT6; D2 — C-7 Human Only)"),
    ("3. Writing the credit record  (compliance step routinely skipped under pressure)",
     "The audit trail fields exist in the system — named approver, case reference, reason code — "
     "but the current process bypasses them. At least one confirmed miss in live records. At machine "
     "speed, this gap would become systematic.  (D1 — WS4 BP-4; Artefact 2 internal note)"),
]

y = Inches(1.6)
colours = [APEX_MID, AMBER, RED_COL]
for (title, detail), col in zip(hotspots, colours):
    # Colour swatch
    swatch = s.shapes.add_shape(1, Inches(0.3), y + Inches(0.04),
                                 Inches(0.12), Inches(0.5))
    swatch.fill.solid()
    swatch.fill.fore_color.rgb = col
    swatch.line.fill.background()

    add_text_box(s, title, Inches(0.55), y, Inches(12.1), Inches(0.42),
                 font_size=15, bold=True, colour=APEX_DARK)
    y += Inches(0.42)
    add_text_box(s, detail, Inches(0.55), y, Inches(12.1), Inches(0.65),
                 font_size=12, colour=GREY_TEXT)
    y += Inches(0.78)

add_footer(s)
add_speaker_notes(s,
    "The first hotspot is the one an agent fixes most directly and most safely — data assembly is not "
    "a judgment call, it's a retrieval task that takes time because the systems aren't integrated. The "
    "second hotspot is the one that requires a policy decision from you before we can build anything. "
    "The third hotspot is the reason the compliance argument for this agent is as strong as the "
    "efficiency argument — we're not just saving time, we're closing a gap that's already creating "
    "audit exposure today.")


# ── SLIDE 7 — Delegation Matrix ───────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "What Can Be Delegated — and What Cannot")

headers7 = ["Task", "Suitable for agent?", "Why"]
rows7 = [
    ["Retrieve invoice, surcharge, and dispute records",            "Agent ✓",       "Structured lookup — no judgment required"],
    ["Send initial acknowledgement to customer",                    "Agent ✓",       "Same message every time; no variation"],
    ["Classify dispute type (fuel surcharge / redelivery / dim. wt.)", "Agent ✓",   "Rule-based once dispute arrives in CRM"],
    ["Verify whether a charge was calculated correctly",            "Agent ✓",       "Arithmetic check against system data"],
    ["Write the completed credit record once approved",             "Agent ✓",       "Structured field population — deterministic"],
    ["Decide the credit amount",                                    "Human only ✗",  "No formal policy exists; judgment required"],
    ["Confirm & sign every credit record before it is written",     "Human — system gate ✗", "Named approver required; enforced by design, not policy"],
]

col_widths7 = [Inches(5.5), Inches(2.2), Inches(5.0)]
tbl7 = s.shapes.add_table(
    len(rows7)+1, 3,
    Inches(0.3), Inches(1.55),
    sum(col_widths7), Inches(4.8)
).table

for ci, (w, h) in enumerate(zip(col_widths7, headers7)):
    tbl7.columns[ci].width = w
    cell = tbl7.cell(0, ci)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = APEX_DARK
    cell.text_frame.paragraphs[0].runs[0].font.bold  = True
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
    cell.text_frame.paragraphs[0].runs[0].font.size  = Pt(12)

for ri, rd in enumerate(rows7, start=1):
    human_row = ri >= 6
    for ci, val in enumerate(rd):
        cell = tbl7.cell(ri, ci)
        cell.text = val
        if human_row:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(11)
        p.runs[0].font.bold = human_row
        p.runs[0].font.color.rgb = RED_COL if human_row else GREY_TEXT

add_footer(s)
add_speaker_notes(s,
    "The right column is not a shortlist of things we couldn't figure out how to automate. It's where "
    "the design draws a deliberate line. The credit amount decision needs a written policy before any "
    "agent can apply it consistently — and that policy doesn't exist yet, which is one of the decisions "
    "I'll come back to. The approval gate is different: that's a system constraint we're building in as "
    "a hard rule. The agent physically cannot write a credit record until a named human has authenticated "
    "and confirmed it in the workflow.")


# ── SLIDE 8 — Opportunity Matrix ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "The Opportunity — Where Volume Meets Complexity")

# Quadrant chart (schematic using shapes)
chart_left = Inches(0.5)
chart_top  = Inches(1.55)
chart_w    = Inches(7.5)
chart_h    = Inches(5.0)

# Axes
ax_h = s.shapes.add_shape(1, chart_left, chart_top + chart_h/2 - Inches(0.02),
                           chart_w, Inches(0.04))
ax_h.fill.solid(); ax_h.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
ax_h.line.fill.background()

ax_v = s.shapes.add_shape(1, chart_left + chart_w/2 - Inches(0.02), chart_top,
                           Inches(0.04), chart_h)
ax_v.fill.solid(); ax_v.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
ax_v.line.fill.background()

# Axis labels
add_text_box(s, "← Low volume          High volume →",
             chart_left, chart_top + chart_h + Inches(0.05), chart_w, Inches(0.3),
             font_size=10, colour=GREY_TEXT, align=PP_ALIGN.CENTER)
add_text_box(s, "High\ncomplexity",
             chart_left - Inches(0.45), chart_top, Inches(0.4), Inches(0.9),
             font_size=9, colour=GREY_TEXT, align=PP_ALIGN.CENTER)
add_text_box(s, "Low\ncomplexity",
             chart_left - Inches(0.45), chart_top + chart_h - Inches(0.9), Inches(0.4), Inches(0.9),
             font_size=9, colour=GREY_TEXT, align=PP_ALIGN.CENTER)

# Plot work streams
ws_positions = [
    # name,   x offset from chart_left, y offset from chart_top, colour, size
    ("WS4\nBilling Disputes\n★ PRIMARY TARGET",
     chart_w * 0.35, chart_h * 0.18, APEX_MID, Inches(1.1), True),
    ("WS1\nDelivery Exceptions",
     chart_w * 0.28, chart_h * 0.3,  GREY_TEXT, Inches(0.9), False),
    ("WS3\nDispatch Adjustments*",
     chart_w * 0.15, chart_h * 0.35, GREY_TEXT, Inches(0.9), False),
    ("WS2\nETA Inquiries",
     chart_w * 0.65, chart_h * 0.75, GREY_TEXT, Inches(0.9), False),
]

for name, xo, yo, col, diam, primary in ws_positions:
    cx = chart_left + xo - diam/2
    cy = chart_top  + yo - diam/2
    dot = s.shapes.add_shape(9, cx, cy, diam, diam)   # oval
    dot.fill.solid()
    dot.fill.fore_color.rgb = APEX_MID if primary else RGBColor(0xAA,0xBB,0xCC)
    dot.line.fill.background()
    add_text_box(s, name, cx - Inches(0.1), cy + diam + Inches(0.04),
                 diam + Inches(0.2), Inches(0.7),
                 font_size=9 if not primary else 10,
                 bold=primary, colour=APEX_DARK if primary else GREY_TEXT,
                 align=PP_ALIGN.CENTER)

# Key metrics panel on the right
panel = s.shapes.add_shape(1, Inches(8.3), Inches(1.55), Inches(4.7), Inches(5.0))
panel.fill.solid(); panel.fill.fore_color.rgb = APEX_LIGHT
panel.line.color.rgb = APEX_MID

metrics = [
    ("Agentic Value Score", "20 / 25  — strongest in portfolio"),
    ("Annual baseline cost", "~£245k / year at current handle time"),
    ("Projected agent cost", "~£70k / year (incl. human review time)"),
    ("Directional saving",   "~£175k / year"),
    ("Build cost estimate",  "~£100k"),
    ("Payback period",       "~7 months"),
]
my = Inches(1.75)
add_text_box(s, "WS4 Billing Disputes", Inches(8.5), my, Inches(4.3), Inches(0.4),
             font_size=14, bold=True, colour=APEX_DARK)
my += Inches(0.5)
for label, val in metrics:
    add_text_box(s, label, Inches(8.5), my, Inches(4.3), Inches(0.28),
                 font_size=11, bold=True, colour=APEX_DARK)
    my += Inches(0.3)
    add_text_box(s, val, Inches(8.5), my, Inches(4.3), Inches(0.28),
                 font_size=12, colour=APEX_MID)
    my += Inches(0.42)

add_footer(s)
add_speaker_notes(s,
    "The chart shows why billing disputes is the right first target. It has the highest combination "
    "of case complexity and daily volume — not so high in volume that it's already been automated, "
    "but high enough in complexity that it genuinely needs an agent rather than a simple automation "
    "script. ETA inquiries are high volume but low complexity — closer to a lookup than an agent. "
    "Dispatch adjustments score similarly to billing disputes, but the dispatch console runs on Citrix "
    "and doesn't have a confirmed API surface, so building there would risk recreating exactly the "
    "brittle integration that broke in 2024.")


# ── SLIDE 9 — The Proposed Agent ─────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "The Proposed Agent — What It Does")

capabilities = [
    ("Receives & assembles evidence",
     "Retrieves invoice data, surcharge records, delivery outcome, and customer account history — "
     "what Sandra currently spends 10–12 minutes retrieving manually.  (D4 T-001 – T-004)"),
    ("Assesses charge validity",
     "Rule-based checks for fuel surcharges, dimensional weight, and redelivery fees. Routes "
     "clear-cut cases with a structured verdict; routes ambiguous cases to a human reviewer "
     "with all evidence assembled and a confidence score.  (D4 T-007)"),
    ("Enforces a complete audit trail on every credit record",
     "Named approver, CRM case reference, approved reason code — every field the current "
     "process routinely leaves blank.  (D4 T-011, FM-3)"),
    ("Flags repeat dispute patterns automatically",
     "A customer with multiple open disputes of the same type triggers senior review rather "
     "than being handled as three separate cases — the Hayes & Sons situation would have been "
     "surfaced on day one.  (D4 T-008, ET-005)"),
]

y = Inches(1.6)
for i, (title, detail) in enumerate(capabilities):
    num_box = s.shapes.add_shape(9, Inches(0.3), y + Inches(0.02),
                                  Inches(0.45), Inches(0.45))
    num_box.fill.solid(); num_box.fill.fore_color.rgb = APEX_MID
    num_box.line.fill.background()
    add_text_box(s, str(i+1), Inches(0.3), y + Inches(0.02), Inches(0.45), Inches(0.45),
                 font_size=14, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, title, Inches(0.9), y, Inches(11.8), Inches(0.38),
                 font_size=14, bold=True, colour=APEX_DARK)
    y += Inches(0.4)
    add_text_box(s, detail, Inches(0.9), y, Inches(11.8), Inches(0.55),
                 font_size=12, colour=GREY_TEXT)
    y += Inches(0.7)

add_footer(s)
add_speaker_notes(s,
    "What this agent replaces is the data assembly and the audit trail enforcement — the parts of "
    "Sandra's job that take the most time but require the least judgment. What it does not replace "
    "is the part that actually requires Sandra: deciding whether a credit is warranted, how much, "
    "and — in the cases that don't fit the standard types — making a call. The target outcome is "
    "that Sandra spends 8 minutes reviewing a pre-assembled case and confirming an amount, rather "
    "than 28 minutes rebuilding the picture from scratch and then writing it up informally.")


# ── SLIDE 10 — Autonomy Boundary ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "Where the Agent Stops — The Autonomy Boundary")

headers10 = ["Tier", "What happens", "Examples"]
rows10 = [
    ["Agent decides alone",
     "Data retrieval, dispute classification, charge calculation check, stale-data flag, "
     "customer acknowledgement, audit record write after approval is confirmed",
     "Retrieve invoice; classify as fuel surcharge dispute; verify surcharge arithmetic; "
     "notify customer once credit is written"],
    ["Agent prepares — named human must approve",
     "Every credit record: agent prepares complete record with all required fields; system holds "
     "it in pending state until named approver authenticates and confirms via authenticated action; "
     "credits above COO-set threshold require senior approver",
     "Sandra reviews case summary, confirms credit amount, her identity is recorded — then and "
     "only then does the system write the record"],
    ["Human only",
     "Credit amount decision (no written policy exists yet); disputes outside the three standard "
     "types; legal or ombudsman referrals; accounts under a payment plan or in collections",
     "Sandra decides the credit amount; agent prepares the paperwork around that decision"],
]

col_w10 = [Inches(2.2), Inches(5.5), Inches(5.0)]
tbl10 = s.shapes.add_table(
    4, 3, Inches(0.3), Inches(1.55), sum(col_w10), Inches(4.9)
).table

tier_colours = [GREEN_COL, AMBER, RED_COL]
tier_bg      = [RGBColor(0xE8,0xF8,0xED), RGBColor(0xFF,0xF5,0xE0), RGBColor(0xFF,0xF0,0xF0)]

for ci, (w, h) in enumerate(zip(col_w10, headers10)):
    tbl10.columns[ci].width = w
    cell = tbl10.cell(0, ci)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = APEX_DARK
    cell.text_frame.paragraphs[0].runs[0].font.bold = True
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)

for ri, (rd, tc, bg) in enumerate(zip(rows10, tier_colours, tier_bg), start=1):
    for ci, val in enumerate(rd):
        cell = tbl10.cell(ri, ci)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(10)
        p.runs[0].font.bold = (ci == 0)
        p.runs[0].font.color.rgb = tc if ci == 0 else GREY_TEXT

# Enforcement callout
add_text_box(s,
    "⚑  The named-approver gate is enforced by the system — not by procedure. "
    "The agent has no write access to the approver field. A daily audit scan detects any bypass within 24 hours.  (D4 §5, FM-5)",
    Inches(0.4), Inches(6.6), Inches(12.3), Inches(0.7),
    font_size=11, colour=APEX_DARK)

add_footer(s)
add_speaker_notes(s,
    "I want to spend a moment on the middle row, because it's the most important one. The reason "
    "Sandra's £170 credit appeared in the Artefact 2 case with no audit log entry is that the "
    "current process relies on people following a procedure under time pressure. That's a "
    "procedure-dependent control — and we've already seen it fail. What we're building is a "
    "system-dependent control: the credit record cannot be written until a human has authenticated. "
    "That's not a rule we're asking your team to follow. It's a technical constraint.")


# ── SLIDE 11 — Integration Readiness ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "Integration Readiness")

headers11 = ["Integration", "Status", "What it means"]
rows11 = [
    ["CRM (Salesforce) — case queue, customer data, delivery records, outbound messaging",
     "AMBER",
     "REST APIs confirmed available. Whether inbound disputes auto-create CRM cases vs. Sandra "
     "creating them manually is not yet confirmed — if manual, the agent's intake trigger depends "
     "on Sandra taking the first step"],
    ["Aurum Billing — credit record write path",
     "RED — BLOCKING",
     "No programmatic write path to the credit ledger confirmed. If none exists, the agent can "
     "prepare credit records but cannot submit them, and the 48-hour manual ticket process remains. "
     "Single largest risk to the build timeline.  (D5 G-1)"],
    ["Credit policy",
     "RED — BLOCKING",
     "No formal written credit policy exists. The agent cannot recommend credit amounts without "
     "explicit approved rules. This is a policy design task — not a technical one — but it must be "
     "completed before deployment with full capability.  (D5 G-2)"],
]

col_w11 = [Inches(4.5), Inches(1.8), Inches(6.4)]
tbl11 = s.shapes.add_table(
    4, 3, Inches(0.3), Inches(1.55), sum(col_w11), Inches(4.0)
).table

status_colours = [AMBER, RED_COL, RED_COL]
status_bg      = [RGBColor(0xFF,0xF5,0xE0), RGBColor(0xFF,0xEC,0xEC), RGBColor(0xFF,0xEC,0xEC)]

for ci, (w, h) in enumerate(zip(col_w11, headers11)):
    tbl11.columns[ci].width = w
    cell = tbl11.cell(0, ci)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = APEX_DARK
    cell.text_frame.paragraphs[0].runs[0].font.bold = True
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)

for ri, (rd, sc, bg) in enumerate(zip(rows11, status_colours, status_bg), start=1):
    for ci, val in enumerate(rd):
        cell = tbl11.cell(ri, ci)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(10)
        p.runs[0].font.bold = (ci == 1)
        p.runs[0].font.color.rgb = sc if ci == 1 else GREY_TEXT

# Callout question
add_text_box(s,
    "Confirmation needed before build can start:\n"
    '"Does Aurum Billing expose any write interface — even a controlled database path or structured '
    'import file — that does not require submitting a manual support ticket?"',
    Inches(0.4), Inches(5.75), Inches(12.3), Inches(1.0),
    font_size=12, colour=APEX_DARK)

add_footer(s)
add_speaker_notes(s,
    "The amber and two reds are not reasons not to build. They're the specific things that need to "
    "be resolved before we can commit to a timeline. The Salesforce amber is likely a configuration "
    "question — your IT team or Salesforce admin should be able to answer it in a day. The two reds "
    "are genuinely blocking: we cannot build the credit recommendation or credit write capability "
    "until they're resolved. The good news is that neither requires a new system to be built. One "
    "is a policy design decision — who owns it, and what are the rules. The other is a question for "
    "your Aurum vendor that we'd like to ask together.")


# ── SLIDE 12 — Five Questions ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "What We Need From You — Five Questions That Determine the Design")

questions = [
    ("1. Credit policy",
     "Does a written credit policy exist anywhere — even a draft or email thread — that documents "
     "the rules for when to credit and how much?\n"
     "If yes: it becomes the agent's rulebook.  If no: someone must write it before build begins."),
    ("2. Sandra's credit-apply steps",
     "When Sandra applies a credit today, what exactly does she do, step by step — which system, "
     "what does she type or click?\n"
     "This single answer tells us whether a programmatic credit write path exists."),
    ("3. Approval authority",
     "Who at Apex is authorised to approve a credit — and is there a threshold above which a "
     "manager must sign off?\n"
     "We need this number to configure the approval routing."),
    ("4. 2024 failure — what would convince you this time is different?",
     "The 2024 RPA broke when Aurum's data format changed without notice. This agent reads from "
     "the same Aurum exports. What would you need to see in the first 30 days?\n"
     "Your answer determines whether we lead deployment with the compliance feature or handle-time reduction."),
    ("5. Sandra's redeployed capacity",
     "If this agent reduces Sandra's time on billing disputes from 28 min to ~8 min per case, "
     "what does she do with the rest of her day — is there already a plan for that capacity?\n"
     "An agent that creates capacity without a redeployment plan generates internal resistance."),
]

y = Inches(1.58)
for q_title, q_detail in questions:
    add_text_box(s, q_title, Inches(0.4), y, Inches(12.4), Inches(0.32),
                 font_size=13, bold=True, colour=APEX_DARK)
    y += Inches(0.33)
    add_text_box(s, q_detail, Inches(0.6), y, Inches(12.2), Inches(0.52),
                 font_size=11, colour=GREY_TEXT)
    y += Inches(0.6)

add_footer(s)
add_speaker_notes(s,
    "I want to be clear about why these are the five questions and not others. We're not asking "
    "about your systems — we have the Aurum exports and the CRM data. We're not asking about "
    "volumes — they're in the scenario. These five are the ones where the answer materially changes "
    "what we build. The credit policy question is the most important: if the answer is that no policy "
    "exists and you're not sure who would own writing one, that changes the build timeline more than "
    "any technical question. I'd like to spend most of our discussion time on questions two and one.")


# ── SLIDE 13 — Discussion ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, APEX_DARK)
add_header_bar(s, "Discussion")

topics = [
    ("1. The compliance gap is already there.",
     "Credits are being applied today without an audit trail — the agent closes this gap "
     "technically, but it will make that gap visible in a way it hasn't been before. "
     "Sandra's process will change.  How does the team hear that?"),
    ("2. The Aurum write path has three realistic options.",
     "Automated email-ticket submission · Agent-assisted manual submit · Direct database write "
     "with Aurum vendor support.  Each has a different risk and cost profile.  "
     "Which fits your risk appetite for the first six months?"),
    ("3. The credit policy is a business decision, not a technical one.",
     "The agent needs explicit written rules before it can recommend credit amounts. "
     "We can help frame what that document needs to contain — but approval and ownership "
     "sit with you and finance.  Is this a pre-build sprint, or does a version already exist?"),
]

y = Inches(2.0)
for title, detail in topics:
    add_text_box(s, title, Inches(0.8), y, Inches(12.0), Inches(0.42),
                 font_size=16, bold=True, colour=WHITE)
    y += Inches(0.44)
    add_text_box(s, detail, Inches(0.8), y, Inches(12.0), Inches(0.7),
                 font_size=13, colour=RGBColor(0xBB, 0xCC, 0xDD))
    y += Inches(0.95)

add_footer(s)
add_speaker_notes(s,
    "These three prompts are the places where the assessment surfaces tensions that can't be "
    "resolved with better data — they need a decision from this room. The first one is about "
    "change management: I'd rather have that conversation now than discover it during rollout. "
    "The second is a genuine fork in the build design with real cost and timeline differences. "
    "The third is the most important: if there is a version of a credit policy somewhere — a "
    "finance memo, a management email, even a whiteboard photo — that's a very different starting "
    "point than writing from scratch.")


# ── SLIDE 14 — Next Steps ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, WHITE)
add_header_bar(s, "Next Steps")

headers14 = ["Action", "Owner", "Dependency", "Target date"]
rows14 = [
    ["Ask Apex IT and the Aurum vendor: does Aurum expose any write interface for credit records "
     "that does not require the manual support ticket?",
     "Apex IT lead + FDE team", "Discovery call scheduled", "[DATE + 2 weeks]"],
    ["Confirm Salesforce configuration: is an Approval Process or Flow available to enforce "
     "the system-level approval gate?",
     "Apex Salesforce admin", "Admin access", "[DATE + 2 weeks]"],
    ["Commission formal credit policy document: explicit rules per dispute type, approved credit "
     "amounts, and named approval threshold",
     "COO + finance lead", "COO decision to proceed", "[DATE + 4 weeks]"],
    ["Identify 150 historical billing dispute cases from CRM archive for confidence calibration; "
     "confirm two senior billing agents available to label them independently",
     "Operations lead", "Credit policy draft complete", "[DATE + 5 weeks]"],
]

col_w14 = [Inches(5.2), Inches(2.0), Inches(2.4), Inches(2.8)]
tbl14 = s.shapes.add_table(
    5, 4, Inches(0.3), Inches(1.55), sum(col_w14), Inches(4.8)
).table

for ci, (w, h) in enumerate(zip(col_w14, headers14)):
    tbl14.columns[ci].width = w
    cell = tbl14.cell(0, ci)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = APEX_DARK
    cell.text_frame.paragraphs[0].runs[0].font.bold = True
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)

row_bg = [APEX_LIGHT, WHITE, APEX_LIGHT, WHITE]
for ri, (rd, bg) in enumerate(zip(rows14, row_bg), start=1):
    for ci, val in enumerate(rd):
        cell = tbl14.cell(ri, ci)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(10)
        p.runs[0].font.color.rgb = GREY_TEXT

add_footer(s)
add_speaker_notes(s,
    "The first two actions are discovery — they can happen in parallel and in under two weeks if "
    "the right people are in the room. The third is the most consequential: it's not a task we "
    "can do for you, but we can provide the template and the framing. The fourth is the pre-deployment "
    "safety requirement — we will not deploy an agent that makes autonomous validity assessments "
    "until we've validated its confidence scores against 150 real cases labelled by your own senior "
    "agents. That's the answer to the question about what's different this time.")


# ── SLIDE 15 — Closing ────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, APEX_DARK)

# Accent stripe
stripe = s.shapes.add_shape(1, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08))
stripe.fill.solid(); stripe.fill.fore_color.rgb = APEX_MID
stripe.line.fill.background()

add_text_box(s, "The Recommendation",
             Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.6),
             font_size=22, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

add_text_box(s,
    "Apex's billing dispute process is the strongest candidate for an AI agent in Customer "
    "Operations — a projected £175k annual saving with a 7-month payback — but two decisions "
    "must be made before build begins: what the credit policy says, and whether Aurum can "
    "accept a programmatic credit write.",
    Inches(1.0), Inches(1.75), Inches(11.3), Inches(1.3),
    font_size=17, colour=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

add_text_box(s, "Contact: [NAME]  |  [EMAIL]",
             Inches(1.0), Inches(3.55), Inches(11.3), Inches(0.4),
             font_size=14, colour=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)
add_text_box(s, "Next meeting: [DATE]",
             Inches(1.0), Inches(3.98), Inches(11.3), Inches(0.4),
             font_size=14, colour=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

add_footer(s)
add_speaker_notes(s,
    "Close by returning to the question the assessment set out to answer: can a targeted agent "
    "reduce dispute handling time, close the compliance gap, and give Apex an auditable trace of "
    "every credit decision without repeating the failures of 2024? The answer is yes — but only "
    "if the two blocking gaps are resolved first. Those are your decisions, and I'd like to leave "
    "this room with clarity on who owns each one and what the timeline is.")


# ── Save ──────────────────────────────────────────────────────────────────────
output_path = r"c:\Users\Benoit_Charrier\FDE Program\FDE-Program\Gate2\Deliverables\Stakeholder_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
