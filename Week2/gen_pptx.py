#!/usr/bin/env python3
"""Generate Stakeholder_Presentation.pptx — Helix Workforce Software ATX Assessment."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colours ────────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1A, 0x3A, 0x5C)
BLUE      = RGBColor(0x2E, 0x75, 0xB6)
LT_BLUE   = RGBColor(0xE8, 0xF0, 0xF7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1A, 0x1A, 0x1A)
GRAY      = RGBColor(0x55, 0x55, 0x55)
TBL_HDR   = RGBColor(0x1F, 0x4E, 0x79)
ALT_ROW   = RGBColor(0xEE, 0xF4, 0xFB)
AMBER     = RGBColor(0xED, 0x7D, 0x31)
RED       = RGBColor(0xC0, 0x00, 0x00)
LT_RED    = RGBColor(0xFF, 0xEB, 0xEB)
LT_GREEN  = RGBColor(0xD6, 0xE8, 0xD4)
DK_GREEN  = RGBColor(0x1E, 0x5C, 0x1E)
LT_STEEL  = RGBColor(0xAD, 0xCC, 0xE8)

SW = Inches(13.333)
SH = Inches(7.5)

# ── Core helpers ───────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s

def box(slide, l, t, w, h, fill=NAVY, line=None):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    lc = line or fill
    sh.line.color.rgb = lc
    return sh

def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)

def fmt(run, text, sz, bold=False, italic=False, color=DARK, font='Calibri'):
    run.text = text
    run.font.name = font; run.font.size = Pt(sz)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color

def p0(tf, text, sz=12, bold=False, italic=False, color=DARK,
       align=PP_ALIGN.LEFT, font='Calibri'):
    """Set first paragraph of a text frame."""
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); fmt(r, text, sz, bold, italic, color, font)
    return p

def ap(tf, text, sz=12, bold=False, italic=False, color=DARK,
       align=PP_ALIGN.LEFT, font='Calibri', spc=0):
    """Add a new paragraph to a text frame."""
    p = tf.add_paragraph(); p.alignment = align
    if spc:
        pPr = p._p.get_or_add_pPr()
        sb  = etree.SubElement(pPr, qn('a:spcBef'))
        sp  = etree.SubElement(sb,  qn('a:spcPts'))
        sp.set('val', str(spc * 100))
    r = p.add_run(); fmt(r, text, sz, bold, italic, color, font)
    return p

def title_bar(s, title):
    b  = box(s, Inches(0), Inches(0), SW, Inches(1.1), fill=NAVY)
    tf = b.text_frame
    tf.margin_left = Inches(0.35); tf.margin_top = Inches(0.2)
    tf.word_wrap   = False
    p0(tf, title, 26, bold=True, color=WHITE)
    box(s, Inches(0), Inches(1.1), SW, Inches(0.04), fill=BLUE, line=BLUE)

def notes(s, text):
    ns = s.notes_slide; tf = ns.notes_text_frame; tf.clear()
    p0(tf, text, 11)

def foot(s, text, t=Inches(7.1)):
    t_ = tb(s, Inches(0.3), t, Inches(12.7), Inches(0.3))
    tf = t_.text_frame
    p0(tf, text, 8, italic=True, color=GRAY)

def cell_fmt(cell, text, bg=WHITE, fc=DARK, sz=10, bold=False,
             align=PP_ALIGN.LEFT, italic=False):
    cell.text = ''
    tf = cell.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.07)
    tf.margin_top  = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); fmt(r, text, sz, bold, italic, fc)
    cell.fill.solid(); cell.fill.fore_color.rgb = bg

def mktable(s, hdrs, rows, l, t, w, h, cws=None,
            hdr_bg=TBL_HDR, alt=True):
    tbl = s.shapes.add_table(len(rows)+1, len(hdrs), l, t, w, h).table
    if cws:
        for i,c in enumerate(cws): tbl.columns[i].width = c
    for j,h_ in enumerate(hdrs):
        cell_fmt(tbl.cell(0,j), h_, bg=hdr_bg, fc=WHITE, sz=11,
                 bold=True, align=PP_ALIGN.CENTER)
    for i,row in enumerate(rows):
        bg = ALT_ROW if (alt and i%2==1) else WHITE
        for j,v in enumerate(row):
            cell_fmt(tbl.cell(i+1,j), v, bg=bg, sz=10)
    return tbl

def callout(s, text, l, t, w, h, bg=LT_BLUE, bc=BLUE, sz=11, bold=False):
    sh = box(s, l, t, w, h, fill=bg, line=bc)
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right  = Inches(0.15)
    tf.margin_top  = tf.margin_bottom = Inches(0.1)
    p0(tf, text, sz, bold=bold, color=DARK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
prs = new_prs()

# ── S1: Title ─────────────────────────────────────────────────────────────────
s1 = blank(prs)
s1.background.fill.solid(); s1.background.fill.fore_color.rgb = NAVY

t_ = tb(s1, Inches(0.7), Inches(1.1), Inches(12), Inches(0.7))
p0(t_.text_frame, 'Helix Workforce Software', 19, color=LT_STEEL)

t_ = tb(s1, Inches(0.7), Inches(1.85), Inches(12), Inches(1.5))
tf = t_.text_frame; tf.word_wrap = True
p0(tf, 'Vendor Contract Clause Review', 38, bold=True, color=WHITE)

t_ = tb(s1, Inches(0.7), Inches(3.5), Inches(12), Inches(0.7))
p0(t_.text_frame, 'Assessment Findings & Proposed Solution', 21, color=LT_STEEL)

box(s1, Inches(0.7), Inches(4.4), Inches(4), Inches(0.05), fill=BLUE, line=BLUE)

t_ = tb(s1, Inches(0.7), Inches(4.55), Inches(7), Inches(0.5))
p0(t_.text_frame, '[Date]  |  [Presenter]', 15, color=LT_STEEL)

notes(s1, "I want to set expectations before we start: this is not a generic AI pitch. We've spent the last several weeks mapping how the Legal & Commercial team actually works — the volume, the bottlenecks, the governance rules — and what we're presenting today is a specific recommendation grounded in that analysis. We'll show you what we found, what we think should be built, and the decisions that only you can make before we can finalise the design.")

# ── S2: Agenda ────────────────────────────────────────────────────────────────
s2 = blank(prs)
title_bar(s2, 'Agenda')

items = [
    ('1.', 'The business problem',  'What we were asked to solve'),
    ('2.', 'What we found',          'Where the work actually goes'),
    ('3.', 'What we recommend',      'The proposed solution'),
    ('4.', 'Decisions we need',      'Open questions that change the design'),
    ('5.', 'Next steps',             'What happens before we can build'),
]
y = Inches(1.4)
for num, sec, desc in items:
    t_ = tb(s2, Inches(0.5), y, Inches(0.5), Inches(0.65))
    p0(t_.text_frame, num, 16, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
    t_ = tb(s2, Inches(1.1), y, Inches(4.0), Inches(0.65))
    p0(t_.text_frame, sec, 16, bold=True, color=NAVY)
    t_ = tb(s2, Inches(5.3), y, Inches(7.8), Inches(0.65))
    p0(t_.text_frame, desc, 15, color=DARK)
    if num != '5.':
        box(s2, Inches(0.5), y+Inches(0.6), Inches(12.5), Inches(0.02),
            fill=LT_BLUE, line=LT_BLUE)
    y += Inches(0.9)

notes(s2, "We'll move through five sections. The first two are findings — we'll earn the recommendation before we make it. The third is the recommendation itself, including where the system stops. The fourth is critical: there are five questions whose answers will materially change what we build. And the fifth is a concrete action list. I'll pause for questions throughout, but I've also reserved time at the end specifically for the decisions we need your input on.")

# ── S3: Business problem ──────────────────────────────────────────────────────
s3 = blank(prs)
title_bar(s3, 'Why We Are Here — The Business Problem')

bullets = [
    ('125 hours of paralegal time per quarter consumed by first-pass clause review — before any negotiation begins',
     '[D3, scenario_context.md]'),
    ('4–6 day review cycle, with CRO pressure to halve it — first-pass classification is the intake gate that everything else waits on',
     '[D1, D3]'),
    ('Active compliance risk: every Data Processing Agreement reviewed today is measured against a playbook that has not incorporated the DPDI Act Q1 changes — 9 months stale',
     '[D4, scenario_context.md]'),
]
y = Inches(1.35)
for text, src in bullets:
    box(s3, Inches(0.35), y+Inches(0.12), Inches(0.38), Inches(0.38), fill=BLUE)
    t_ = tb(s3, Inches(0.88), y, Inches(11.5), Inches(0.75))
    tf = t_.text_frame; tf.word_wrap = True
    p0(tf, text, 14, color=DARK)
    t_ = tb(s3, Inches(0.88), y+Inches(0.72), Inches(5), Inches(0.28))
    p0(t_.text_frame, src, 9, italic=True, color=GRAY)
    y += Inches(1.3)

box(s3, Inches(0.35), y, Inches(12.63), Inches(0.04), fill=BLUE, line=BLUE)
callout(s3,
    'Assessment question: Can first-pass clause classification be reliably automated, '
    'with the named-lawyer sign-off requirement preserved and auditable — not assumed?',
    Inches(0.35), y+Inches(0.12), Inches(12.63), Inches(0.75),
    bg=LT_BLUE, bc=BLUE, sz=13)

notes(s3, "The 125-hour figure is worth pausing on. That's the equivalent of roughly three full working weeks every quarter that Tom spends reading contracts and comparing clause text against the playbook — before a single redline is drafted or a single lawyer is consulted. The compliance risk on the DPA clause exists in your process today, agent or no agent. We'll come back to it, because it has direct implications for what we can build and when.")

# ── S4: Process flow ──────────────────────────────────────────────────────────
s4 = blank(prs)
title_bar(s4, 'How the Work Actually Flows Today')

flow = (
    "Vendor contract (.docx, ~25 pages)\n"
    "    |  via Outlook email\n"
    "    v\n"
    "WS1: First-pass clause classification        300 contracts/quarter  x  ~25 min each\n"
    "    |\n"
    "    +-- 70% standard   -->  Accept as-is                          (~210 contracts)\n"
    "    +-- 20% deviation  -->  WS2: Paralegal redlining      (~60 contracts, ~45 min)\n"
    "    +-- 10% escalation -->  WS3: Senior lawyer review     (~30 contracts, ~90 min)\n"
    "                                       |\n"
    "                                       v\n"
    "                 WS4: Counteroffer package + sign-off    (~90 contracts, ~30 min)\n"
    "                                       |\n"
    "                                       v\n"
    "                                 Vendor dispatch"
)
t_ = tb(s4, Inches(0.5), Inches(1.28), Inches(12.3), Inches(4.9))
tf = t_.text_frame; tf.word_wrap = False
p0(tf, flow, 11.5, color=DARK, font='Courier New')

callout(s4,
    'Most cognitive effort: WS1 — the classification step that determines every contract\'s path  [D1, scenario_context.md]',
    Inches(0.35), Inches(6.3), Inches(12.63), Inches(0.65),
    bg=LT_BLUE, bc=BLUE, sz=12, bold=True)

notes(s4, "WS1 is not just one of four work streams — it's the gate. Nothing moves to WS2, WS3, or WS4 until WS1 is done. At 25 minutes per contract across 300 contracts, that's where the time goes and where the turnaround delay originates. The 70/20/10 split is the other key fact: the majority of contracts are standard. The challenge is that you have to read every one of them to know which category they're in. That's the problem we're trying to solve.")

# ── S5: Cognitive hotspots ────────────────────────────────────────────────────
s5 = blank(prs)
title_bar(s5, 'Where Time Goes — Three Judgment Calls That Slow Everything Down')

hotspots = [
    ('1  |  Deciding whether a clause deviation is serious enough to escalate',
     'Tom compares extracted clause language against playbook positions across 7 clause types. For qualitative clauses — IP ownership, indemnity scope — this requires judging commercial intent, not just matching wording. Judgment call that varies case by case.',
     'Automatable with structured oversight  [D1]'),
    ('2  |  Locating the right clause when vendor documents use non-standard headings',
     'A clause titled "Commercial Exposure" may contain liability cap language. Tom recognises this through experience. Without a trained pattern library, an agent misses it entirely.',
     'Requires a curated knowledge base before automation is viable  [D1, D5]'),
    ('3  |  Assessing DPA clauses against a compliance baseline that is not current',
     "Tom reviews Data Processing Agreements against a playbook he knows is 9 months stale. He cannot flag DPDI Act gaps reliably without Amelia's completed update.",
     'Every DPA review today carries latent compliance risk — agent or no agent  [D1, D4]'),
]
y = Inches(1.28)
for hdr, body, status in hotspots:
    box(s5, Inches(0.35), y, Inches(12.63), Inches(0.44), fill=TBL_HDR)
    t_ = tb(s5, Inches(0.5), y+Inches(0.06), Inches(12.3), Inches(0.35))
    p0(t_.text_frame, hdr, 12, bold=True, color=WHITE)
    t_ = tb(s5, Inches(0.5), y+Inches(0.5), Inches(12.1), Inches(0.72))
    tf = t_.text_frame; tf.word_wrap = True
    p0(tf, body, 11, color=DARK)
    t_ = tb(s5, Inches(0.5), y+Inches(1.22), Inches(12.1), Inches(0.26))
    p0(t_.text_frame, status, 10, italic=True, color=BLUE)
    y += Inches(1.85)

notes(s5, "The first hotspot is the one we can automate most readily — it follows a pattern, even if that pattern isn't perfectly deterministic. The second is a knowledge gap we can fill by training the system on historical contract structures. The third is the one that keeps me up at night: it's not an AI problem, it's a data quality problem. An agent trained on the current playbook will produce the same compliance-risk classifications that Tom produces today — at scale. That's why the playbook update isn't optional background work. It's a deployment gate.")

# ── S6: Delegation ────────────────────────────────────────────────────────────
s6 = blank(prs)
title_bar(s6, 'What Can Be Delegated to an Agent — and What Cannot')

h6  = ['Agent-suitable', 'Human-anchored']
r6  = [
    ['Contract intake, routing, and case record creation — structured task, no judgment required',
     'Named-lawyer counteroffer sign-off  [GC rule]\nNo counteroffer may leave Legal\'s queue without a named lawyer\'s approval recorded in Ironclad. Non-negotiable hard stop.  [D2, D4]'],
    ['Clause text extraction from vendor documents — pattern recognition against known structures',
     'Senior lawyer review of unusual clauses outside the 7 playbook categories — no policy position exists for these'],
    ['Standard clause comparison — matching extracted text against playbook positions for the 70% of contracts where deviation is clear',
     'DPA clause review while the DPDI Act playbook update is outstanding — mandatory human review on every DPA case'],
    ['Exception flagging — when deviation magnitude exceeds defined thresholds, routing is deterministic',
     'Redline drafting for qualitative clause types (IP, indemnity) — synthesis judgment; no templatable output'],
    ['Counteroffer package assembly — compiling approved redlines into a structured sign-off package', ''],
]
cw6 = [Inches(6.1), Inches(6.1)]
tbl6 = mktable(s6, h6, r6, Inches(0.35), Inches(1.25), Inches(12.63), Inches(5.95), cws=cw6)

# Highlight GC rule cell
gc = tbl6.cell(1, 1)
gc.fill.solid(); gc.fill.fore_color.rgb = LT_RED
for para in gc.text_frame.paragraphs:
    for run in para.runs:
        run.font.bold = True; run.font.color.rgb = RED

notes(s6, "The right column is not a negotiating position — it reflects the governance constraint you've already built into the process over 12 years. The GC hard rule is listed first because the architecture is designed around it: the system is built so the agent literally cannot dispatch a counteroffer without a named lawyer's approval token being present in the case record. It's not a warning — it's a hard architectural stop.")

# ── S7: Volume x Value ────────────────────────────────────────────────────────
s7 = blank(prs)
title_bar(s7, 'The Opportunity — Where Volume Meets Complexity')

h7  = ['', 'Low complexity\n(routine, pattern-based)', 'High complexity\n(judgment-intensive, varies case-by-case)']
r7  = [
    ['High\nvolume', '—', 'WS1: Clause classification\nPrimary target  |  Score: 12 / 25'],
    ['Low\nvolume',  '—', 'WS2: Redlining (Score 5)  |  WS3: Escalations (Score 5)  |  WS4: Counteroffers (Score 6)'],
]
cw7 = [Inches(1.5), Inches(4.3), Inches(6.8)]
tbl7 = mktable(s7, h7, r7, Inches(0.35), Inches(1.25), Inches(12.63), Inches(3.3), cws=cw7)

# Highlight WS1 cell
ws1 = tbl7.cell(1, 2)
ws1.fill.solid(); ws1.fill.fore_color.rgb = LT_GREEN
for para in ws1.text_frame.paragraphs:
    for run in para.runs:
        run.font.bold = True; run.font.color.rgb = DK_GREEN

t_ = tb(s7, Inches(0.35), Inches(4.65), Inches(12.63), Inches(0.28))
p0(t_.text_frame,
   'Score = Volume x Complexity on a 1-25 scale.  >=15: strong candidate.  8-14: consider, validate with economics.  <8: do not automate.  [D3]',
   9, italic=True, color=GRAY)

callout(s7,
    'Directional finding: ~£43,000 annual saving (WS1 + counteroffer package prep)  |  '
    '~£60,000 estimated build cost  |  ~17-month payback\n'
    'All figures are estimates based on UK rate assumptions — to be validated.  [D3 §8]',
    Inches(0.35), Inches(5.1), Inches(12.63), Inches(1.0),
    bg=LT_BLUE, bc=BLUE, sz=12)

notes(s7, "WS2, WS3, and WS4 score high on complexity but low on volume — they're difficult cases, but there aren't enough of them to justify a standalone automation. WS1 is the only work stream where the volume justifies the build cost and where the judgment pattern is consistent enough to replicate. A score of 12 out of 25 means 'consider building, and validate the conditions before you commit' — it's a strong conditional case. The 17-month payback is directional and should be confirmed against actual rates before signing off on a build budget.")

# ── S8: The agent ─────────────────────────────────────────────────────────────
s8 = blank(prs)
title_bar(s8, 'The Proposed Solution — What the Agent Does')

badge = box(s8, Inches(0.35), Inches(1.25), Inches(4.8), Inches(0.52), fill=TBL_HDR)
t_ = tb(s8, Inches(0.5), Inches(1.31), Inches(4.6), Inches(0.4))
p0(t_.text_frame, 'Agent: Clause Classification Agent (CCA)  [D4]', 12, bold=True, color=WHITE)

verbs = [
    ('Receives',  "every inbound vendor contract via Ironclad, reads the full document, and extracts the text for each of the 7 clause types the playbook covers"),
    ('Compares',  "each extracted clause against the current Helix playbook position and assigns a classification — standard, deviation, or escalation-required — with a certainty level on each decision"),
    ('Routes',    "standard-path contracts (all 7 clauses within playbook tolerances and certainty >= 85%) for acceptance without Tom's full review — cutting WS1 time from 25 minutes to under 5 minutes for those ~210 contracts per quarter"),
    ('Prepares',  "a structured deviation summary for the ~30% of contracts where clauses fall outside playbook tolerances — Tom reviews the agent's findings, not the full document"),
]
y = Inches(2.02)
for verb, desc in verbs:
    vb = box(s8, Inches(0.35), y, Inches(1.45), Inches(0.82), fill=BLUE)
    t_ = tb(s8, Inches(0.4), y+Inches(0.22), Inches(1.35), Inches(0.4))
    p0(t_.text_frame, verb, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    t_ = tb(s8, Inches(1.95), y+Inches(0.1), Inches(11.0), Inches(0.7))
    tf = t_.text_frame; tf.word_wrap = True
    p0(tf, desc, 12, color=DARK)
    y += Inches(1.02)

box(s8, Inches(0.35), y+Inches(0.05), Inches(12.63), Inches(0.03), fill=LT_BLUE, line=LT_BLUE)
t_ = tb(s8, Inches(0.35), y+Inches(0.12), Inches(12.63), Inches(0.5))
tf = t_.text_frame; tf.word_wrap = True
p0(tf, 'Replaces: 125 hours/quarter of end-to-end contract reading for initial clause comparisons  |  '
       'Produces: per-clause classification report, routing decision, and structured input to the counteroffer pipeline  [D4]',
   10, italic=True, color=GRAY)

notes(s8, "I want to be concrete about what 'reduces Tom's time' actually means. For the 70% standard contracts, Tom currently reads every page to confirm nothing deviates. With this agent, he receives a notification: all 7 clauses within playbook tolerance — no action required. For the 30% with deviations, Tom currently does the comparison himself. With this agent, he receives a structured summary showing the clause, the playbook position, how far they diverge, and the proposed routing. He makes the decision; he doesn't redo the comparison. That's the value proposition.")

# ── S9: Autonomy boundary ─────────────────────────────────────────────────────
s9 = blank(prs)
title_bar(s9, 'Where the Agent Stops — The Approval Boundary')

h9 = ['Agent decides alone', 'Agent proposes,\nhuman approves', 'Human only']
r9 = [
    ['Contract intake, case record creation, clause extraction from document',
     'Any contract with one or more deviating clauses — Tom approves the routing before the case moves forward',
     'Named-lawyer counteroffer sign-off  [GC rule]\nApproval token must be recorded in Ironclad by the lawyer before any counteroffer proceeds. Enforced by system design, not policy.'],
    ['Standard-path classification with certainty >= 85% across all 7 clause types',
     'All DPA clause assessments while the DPDI Act playbook update remains outstanding',
     'Senior lawyer review of clause types outside the 7 playbook categories — no automation path exists'],
    ["Routing standard contracts to the 'accept' queue — Tom notified, no review required",
     "Any classification where the agent's certainty falls below 85% — Tom reviews before the decision commits",
     '—'],
]
cw9 = [Inches(4.0), Inches(4.35), Inches(4.28)]
tbl9 = mktable(s9, h9, r9, Inches(0.35), Inches(1.25), Inches(12.63), Inches(5.1), cws=cw9)

gc9 = tbl9.cell(1, 2)
gc9.fill.solid(); gc9.fill.fore_color.rgb = LT_RED
for para in gc9.text_frame.paragraphs:
    for run in para.runs:
        run.font.bold = True; run.font.color.rgb = RED

callout(s9,
    "The approval gate is enforced at database level: the agent's credentials cannot write to the sign-off field. "
    "No instruction — from a downstream system or a human operator — can override this constraint.  [Deliverables/CLAUDE.md]",
    Inches(0.35), Inches(6.5), Inches(12.63), Inches(0.65),
    bg=LT_RED, bc=RED, sz=11)

notes(s9, "The phrase 'enforced by design, not by policy' is the critical distinction. Policy says 'lawyers must sign off.' Design means the database field that records the sign-off cannot be written by the agent — the agent's API credentials are denied write access to that field entirely. A counteroffer dispatch that depends on that field being non-empty cannot proceed without a lawyer physically taking an action in Ironclad. This is not a trust-the-AI question. It's an architecture question, and the answer is: the system cannot bypass this gate.")

# ── S10: Integration readiness ────────────────────────────────────────────────
s10 = blank(prs)
title_bar(s10, 'Integration Readiness')

h10 = ['Integration', 'Status', 'What it means']
r10 = [
    ['Ironclad CLM (most critical)\nREST API confirmed; per-clause classification fields require custom configuration',
     'AMBER',
     'API confirmed. ~35 custom fields needed across 7 clause types. Must be scoped with the Ironclad admin before build begins.  [D5 Gap G-4]'],
    ['HITL review channel\nNo confirmed mechanism for routing flagged contracts to Tom',
     'RED — BLOCKING',
     'Agent can classify but cannot deliver flagged results to Tom. Entire oversight workflow blocked until resolved. Options: Ironclad-native workflow or Outlook shared inbox.  [D5 Gap G-1]'],
    ['DPDI Act regulatory reference\nDocument not yet produced; playbook 9 months stale',
     'RED — DEPLOYMENT GATE',
     "Agent cannot classify DPA clauses against DPDI Act Q1 changes until Amelia's update is completed and loaded. All DPA clauses are mandatory human review until resolved.  [D5 Gap G-2, D4 §8]"],
]
cw10 = [Inches(4.2), Inches(1.9), Inches(6.5)]
tbl10 = mktable(s10, h10, r10, Inches(0.35), Inches(1.25), Inches(12.63), Inches(4.1), cws=cw10)

status_colors = [AMBER, RED, RED]
for i, clr in enumerate(status_colors):
    sc = tbl10.cell(i+1, 1)
    sc.fill.solid(); sc.fill.fore_color.rgb = WHITE
    for para in sc.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = True; run.font.color.rgb = clr; run.font.size = Pt(10)

callout(s10,
    'Confirmation needed before build starts: Can Ironclad be configured with '
    '~35 custom per-clause classification fields? Ironclad admin must confirm schema expandability.  [D5]',
    Inches(0.35), Inches(5.5), Inches(12.63), Inches(0.65),
    bg=LT_BLUE, bc=BLUE, sz=11)

notes(s10, "I want to be honest here: we have two blocking gaps and one deployment gate. This is not 'we're ready to build.' The HITL channel is a design decision as much as a technical one — Tom needs to tell us how he wants to receive flagged contracts. The Ironclad field schema is a confirmation — the Ironclad admin either supports custom fields or they don't. The DPDI Act update is on your critical path, Amelia, and it currently has no owner and no date.")

# ── S11: Questions ────────────────────────────────────────────────────────────
s11 = blank(prs)
title_bar(s11, 'What We Need From You — Five Questions That Change the Design')

qs = [
    ('Q1  How does Tom currently decide whether a clause deviation is serious enough to escalate?',
     'Does he compare wording against specific playbook language, or judge whether the commercial intent is equivalent even if the wording differs?',
     'The answer determines how we design the comparison logic — and how accurate the agent can realistically be at launch.  [D6 Q4]'),
    ('Q2  When a lawyer signs off on a counteroffer today, where does that approval live?',
     'Is it a recorded action in Ironclad, an email, a verbal confirmation — or is it not formally recorded at all?',
     'If sign-off currently happens outside Ironclad, the approval gate architecture requires a process change before we can build it.  [D6 Q7]'),
    ('Q3  Can any of the three commercial lawyers sign off on any clause type, or does authority vary by clause?',
     'For example, is Amelia the only one who can approve a DPA deviation?',
     'If authority varies by clause type, the system must route each case to the correct lawyer — a different routing design.  [D6 Q8]'),
    ('Q4  When a lawyer approves, does one sign-off cover the whole counteroffer package, or does each deviated clause need its own approval?',
     'Contract-level approval is one field in the system. Clause-level approval is a much more complex approval workflow.',
     '[D6 Q9]'),
    ('Q5  The DPDI Act playbook update has been in discussion since March. Is there a named owner and a committed completion date?',
     'This is a deployment gate: the agent cannot classify DPA clauses reliably without it.',
     '[D6 Q16]'),
]
y = Inches(1.3)
for q_title, q_body, q_why in qs:
    box(s11, Inches(0.35), y+Inches(0.04), Inches(0.34), Inches(0.34), fill=BLUE)
    t_ = tb(s11, Inches(0.82), y, Inches(12.0), Inches(0.38))
    tf = t_.text_frame; tf.word_wrap = True
    p0(tf, q_title, 11, bold=True, color=NAVY)
    t_ = tb(s11, Inches(0.82), y+Inches(0.37), Inches(12.0), Inches(0.33))
    tf = t_.text_frame; tf.word_wrap = True
    p0(tf, q_body, 10.5, color=DARK)
    t_ = tb(s11, Inches(0.82), y+Inches(0.68), Inches(12.0), Inches(0.25))
    p0(t_.text_frame, q_why, 9.5, italic=True, color=BLUE)
    y += Inches(1.15)

notes(s11, "These five questions are ones we genuinely cannot answer from what we already know about your process — and each one has a direct impact on what we build. Q1 affects the confidence threshold design. Q2, Q3, and Q4 together determine the entire approval token architecture: if sign-off is currently informal or outside Ironclad, we need to address that as a process gap before we can enforce it in the system. Q5 is the one I'd most like to resolve today, because it's blocking a specific part of the build and it's currently unowned.")

# ── S12: Discussion ───────────────────────────────────────────────────────────
s12 = blank(prs)
title_bar(s12, 'Discussion')

t_ = tb(s12, Inches(0.35), Inches(1.25), Inches(12.63), Inches(0.38))
p0(t_.text_frame, 'Three open questions for your reaction:', 14, bold=True, color=NAVY)

disc = [
    ('On autonomous operation',
     "We've proposed that the agent handles the 70% of standard-path contracts fully autonomously — Tom receives a notification but doesn't review the classification unless he chooses to. Does that level of autonomy feel appropriate from day one, or would you want Tom to spot-check every agent classification for an initial period before the autonomous path goes live?"),
    ('On the governance gate',
     "The system is designed so that no counteroffer can leave Legal's queue without a named lawyer's approval token recorded in Ironclad — that's an architectural rule, not a workflow suggestion. Does sign-off currently happen in Ironclad as a field action, or would formalising it there require a change to how the team currently works?"),
    ('On the deployment timeline',
     "The DPDI Act playbook update is on the critical path for the agent handling DPA clauses. Who is the right person to own that update, and is a completion date achievable within the deployment planning window?"),
]
y = Inches(1.82)
for topic, prompt in disc:
    box(s12, Inches(0.35), y, Inches(2.9), Inches(0.4), fill=BLUE)
    t_ = tb(s12, Inches(0.45), y+Inches(0.06), Inches(2.8), Inches(0.32))
    p0(t_.text_frame, topic, 11, bold=True, color=WHITE)
    t_ = tb(s12, Inches(3.4), y, Inches(9.58), Inches(0.9))
    tf = t_.text_frame; tf.word_wrap = True
    p0(tf, prompt, 12, color=DARK)
    y += Inches(1.62)

notes(s12, "I've deliberately chosen these three because they represent tensions in the design that we can't resolve without your input. The first is about trust calibration — how much autonomous operation is Amelia comfortable with before there's a measured accuracy track record? The second is about whether the governance gate we've designed maps onto current operational reality or requires a process change. The third is the practical blocker: without an owner and a date on the DPDI Act update, the DPA portion of the agent is indefinitely stalled.")

# ── S13: Next steps ───────────────────────────────────────────────────────────
s13 = blank(prs)
title_bar(s13, 'Next Steps')

h13 = ['Action', 'Owner', 'Dependency', 'Target date']
r13 = [
    ['Confirm Ironclad custom field support: ~35 per-clause classification fields across 7 clause types  [D5 G-4 — BLOCKING]',
     'Ironclad admin', 'Ironclad admin access; field schema review', '[Placeholder]'],
    ['Name owner and agree completion date for DPDI Act playbook update — deployment gate for DPA clause handling  [D5 G-2, D6 Q16 — BLOCKING]',
     'Amelia (GC)', 'Internal decision; no external dependency', '[Placeholder]'],
    ['Confirm HITL review channel: how should the agent deliver flagged contracts to Tom?  [D5 G-1 — BLOCKING]',
     'FDE team + Tom', "Tom's workflow preference; Ironclad or Outlook API access", '[Placeholder]'],
    ['Provide 20-30 historical vendor contracts from Ironclad for clause heading pattern analysis  [D5 G-3]',
     'Legal team', 'Access to Ironclad case archive', '[Placeholder]'],
]
cw13 = [Inches(5.7), Inches(1.85), Inches(3.1), Inches(1.85)]
mktable(s13, h13, r13, Inches(0.35), Inches(1.25), Inches(12.63), Inches(5.45), cws=cw13)

foot(s13, 'Actions 1, 2, and 3 are blocking — build specification cannot be finalised without them. Action 4 improves accuracy before go-live but does not block the build.  [D5 G-1, G-2, G-3, G-4; D6 Q7, Q16]')

notes(s13, "These four actions are all that stands between where we are today and a finalised build specification. Two of them — the Ironclad field confirmation and the HITL channel decision — are technical questions with a one-week resolution path. The DPDI update requires Amelia's decision about ownership and timeline; without it, we're committing to a build that excludes DPA clause automation indefinitely. The historical contracts request is the fastest path to improving clause location accuracy before launch.")

# ── S14: Closing ──────────────────────────────────────────────────────────────
s14 = blank(prs)
s14.background.fill.solid(); s14.background.fill.fore_color.rgb = NAVY

box(s14, Inches(0.9), Inches(1.2), Inches(11.5), Inches(2.2), fill=TBL_HDR)
t_ = tb(s14, Inches(1.1), Inches(1.35), Inches(11.13), Inches(2.0))
tf = t_.text_frame; tf.word_wrap = True
p0(tf, 'First-pass clause classification is a strong candidate for automation — with a projected '
       '17-month payback, 125 hours of paralegal time recovered per quarter, and a design that '
       'preserves the named-lawyer sign-off requirement by architectural constraint, not policy assumption.',
   17, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

box(s14, Inches(3.5), Inches(3.75), Inches(6.33), Inches(0.05), fill=BLUE, line=BLUE)

t_ = tb(s14, Inches(0.9), Inches(3.95), Inches(11.5), Inches(0.6))
p0(t_.text_frame, '[Contact details]  |  [Next meeting / follow-up date]',
   15, color=LT_STEEL, align=PP_ALIGN.CENTER)

t_ = tb(s14, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.7))
p0(t_.text_frame, 'Thank you', 30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

notes(s14, "To close: the case for building this agent is solid, conditional. The volume is there, the pattern is there, and the economics close. What makes this different from a generic AI implementation is the governance design — the sign-off gate is not a prompt or a policy, it's a system-level constraint that the agent cannot override. What we need from you today — or within the next week — are the five answers from the questions slide. With those, we can finalise the specification and move to build.")

# ── Save ──────────────────────────────────────────────────────────────────────
out = r"C:\Users\Benoit_Charrier\FDE Program\FDE-Program\Week2\Deliverables\Stakeholder_Presentation.pptx"
prs.save(out)
print(f"Saved {len(prs.slides)} slides -> {out}")
